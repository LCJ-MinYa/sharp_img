import os
from pptx import Presentation
from config import Config

# 指定包含PPT文件的文件夹
folder_path = "C:/Users/Administrator/Desktop/测试/ppt1010"


def del_slide(presentation, slide_index):
    """
    删除某一张幻灯片并返回该幻灯片对象
    :param presentation: Presentation 对象
    :param slide_index: 索引
    :return: 被删除的幻灯片对象
    """
    slides = list(presentation.slides._sldIdLst)
    slide_to_delete = slides[slide_index]
    presentation.slides._sldIdLst.remove(slide_to_delete)
    return slide_to_delete  # 返回被删除的幻灯片对象


def delete_last_slide_from_ppt(ppt_path):
    """
    从PPT中删除最后一张幻灯片
    :param ppt_path: PPT文件路径
    """
    prs = Presentation(ppt_path)
    last_slide_index = len(prs.slides) - 1
    if last_slide_index >= 0:
        del_slide(prs, last_slide_index)
        prs.save(ppt_path)
        print(f"已删除最后一张幻灯片： {ppt_path}")
    else:
        print(f"pptx文件中没有幻灯片： {ppt_path}")


def batch_delete_last_slides(directory):
    """
    批量删除每个PPT的最后一张幻灯片
    :param directory: 文件夹路径
    """
    for filename in os.listdir(directory):
        if filename.endswith(".pptx") or filename.endswith(".ppt"):
            ppt_path = os.path.join(directory, filename)
            delete_last_slide_from_ppt(ppt_path)


# 调用批量删除函数
batch_delete_last_slides(Config.get_latest_folder(Config.WORK_PATH))
